from pptx import Presentation
import io
import logging
from typing import Dict, Any
from datetime import datetime

logger = logging.getLogger(__name__)

class PowerPointExtractor:
    @staticmethod
    def extract_content(file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract text content and metadata from PowerPoint file
        """
        try:
            start_time = datetime.now()
            logger.info(f"Starting PowerPoint extraction for file: {filename}")
            
            # Create presentation object from bytes
            presentation = Presentation(io.BytesIO(file_content))
            
            # Extract basic metadata
            slide_count = len(presentation.slides)
            logger.debug(f"Found {slide_count} slides in {filename}")
            
            # Extract text from slides
            slides_content = []
            all_text = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        all_text.append(shape.text.strip())
                
                slides_content.append({
                    "slide_number": i + 1,
                    "text_content": slide_text,
                    "combined_text": " ".join(slide_text),
                    "shape_count": len(slide.shapes)
                })
            
            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            
            # Prepare extracted content
            extracted_content = {
                "filename": filename,
                "file_size_bytes": len(file_content),
                "slide_count": slide_count,
                "slides": slides_content,
                "all_text_combined": " ".join(all_text),
                "word_count": len(" ".join(all_text).split()),
                "character_count": len(" ".join(all_text)),
                "processing_time_ms": processing_time,
                "metadata": {
                    "total_slides": slide_count,
                    "has_content": len(all_text) > 0,
                    "extraction_timestamp": datetime.now().isoformat(),
                    "extractor_version": "python-pptx-0.6.21"
                }
            }
            
            logger.info(f"Successfully extracted content from {filename}: {slide_count} slides, {len(all_text)} text elements")
            return extracted_content
            
        except Exception as e:
            logger.error(f"Failed to extract PowerPoint content from {filename}: {str(e)}")
            raise ValueError(f"Failed to extract PowerPoint content: {str(e)}")
