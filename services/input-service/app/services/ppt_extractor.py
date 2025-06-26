from pptx import Presentation
import io
import logging
import subprocess
import tempfile
import os
from typing import Dict, Any
from datetime import datetime
from pathlib import Path

logger = logging.getLogger(__name__)

class PowerPointExtractor:
    @staticmethod
    def extract_content(file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract text content and metadata from PowerPoint file (.ppt or .pptx)
        """
        try:
            start_time = datetime.now()
            logger.info(f"Starting PowerPoint extraction for file: {filename}")
            
            file_extension = Path(filename).suffix.lower()
            
            if file_extension == ".pptx":
                return PowerPointExtractor._extract_pptx_content(file_content, filename, start_time)
            elif file_extension == ".ppt":
                return PowerPointExtractor._extract_ppt_content(file_content, filename, start_time)
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")
                
        except Exception as e:
            logger.error(f"Failed to extract PowerPoint content from {filename}: {str(e)}")
            raise ValueError(f"Failed to extract PowerPoint content: {str(e)}")
    
    @staticmethod
    def _extract_pptx_content(file_content: bytes, filename: str, start_time: datetime) -> Dict[str, Any]:
        """Extract content from .pptx files using python-pptx"""
        try:
            # Create presentation object from bytes
            presentation = Presentation(io.BytesIO(file_content))
            
            # Extract basic metadata
            slide_count = len(presentation.slides)
            logger.debug(f"Found {slide_count} slides in {filename}")
            
            # Extract text from slides
            slides_content = []
            all_text = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        all_text.append(shape.text.strip())
                
                slides_content.append({
                    "slide_number": i + 1,
                    "text_content": slide_text,
                    "combined_text": " ".join(slide_text),
                    "shape_count": len(slide.shapes)
                })
            
            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            
            # Prepare extracted content
            extracted_content = {
                "filename": filename,
                "file_size_bytes": len(file_content),
                "slide_count": slide_count,
                "slides": slides_content,
                "all_text_combined": " ".join(all_text),
                "word_count": len(" ".join(all_text).split()),
                "character_count": len(" ".join(all_text)),
                "processing_time_ms": processing_time,
                "metadata": {
                    "total_slides": slide_count,
                    "has_content": len(all_text) > 0,
                    "extraction_timestamp": datetime.now().isoformat(),
                    "extractor_version": "python-pptx-0.6.21",
                    "file_format": "pptx"
                }
            }
            
            logger.info(f"Successfully extracted content from {filename}: {slide_count} slides, {len(all_text)} text elements")
            return extracted_content
            
        except Exception as e:
            logger.error(f"Failed to extract PPTX content: {str(e)}")
            raise
    
    @staticmethod
    def _extract_ppt_content(file_content: bytes, filename: str, start_time: datetime) -> Dict[str, Any]:
        """Extract content from .ppt files using LibreOffice conversion"""
        temp_dir = None
        try:
            # Create temporary directory
            temp_dir = tempfile.mkdtemp()
            
            # Write .ppt file to temp directory
            ppt_path = os.path.join(temp_dir, filename)
            with open(ppt_path, 'wb') as f:
                f.write(file_content)
            
            # Convert .ppt to .pptx using LibreOffice
            pptx_path = PowerPointExtractor._convert_ppt_to_pptx(ppt_path, temp_dir)
            
            if pptx_path and os.path.exists(pptx_path):
                # Read the converted .pptx file
                with open(pptx_path, 'rb') as f:
                    pptx_content = f.read()
                
                # Extract content using pptx method
                result = PowerPointExtractor._extract_pptx_content(pptx_content, filename, start_time)
                result["metadata"]["file_format"] = "ppt"
                result["metadata"]["extractor_version"] = "libreoffice + python-pptx"
                return result
            else:
                # Fallback: Try basic text extraction
                return PowerPointExtractor._extract_ppt_fallback(file_content, filename, start_time)
                
        except Exception as e:
            logger.warning(f"LibreOffice conversion failed for {filename}: {str(e)}")
            # Fallback to basic extraction
            return PowerPointExtractor._extract_ppt_fallback(file_content, filename, start_time)
        finally:
            # Cleanup temporary files
            if temp_dir and os.path.exists(temp_dir):
                try:
                    import shutil
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    logger.warning(f"Failed to cleanup temp directory: {str(e)}")
    
    @staticmethod
    def _convert_ppt_to_pptx(ppt_path: str, output_dir: str) -> str:
        """Convert .ppt to .pptx using LibreOffice"""
        try:
            # Check if LibreOffice is available
            result = subprocess.run(['libreoffice', '--version'], 
                                  capture_output=True, text=True, timeout=10)
            if result.returncode != 0:
                raise Exception("LibreOffice not available")
            
            # Convert using LibreOffice headless mode
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pptx',
                '--outdir', output_dir,
                ppt_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # Find the converted file
                base_name = os.path.splitext(os.path.basename(ppt_path))[0]
                pptx_path = os.path.join(output_dir, f"{base_name}.pptx")
                if os.path.exists(pptx_path):
                    logger.info(f"Successfully converted {ppt_path} to {pptx_path}")
                    return pptx_path
            
            logger.warning(f"LibreOffice conversion failed: {result.stderr}")
            return None
            
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice conversion timed out")
            return None
        except Exception as e:
            logger.warning(f"LibreOffice conversion error: {str(e)}")
            return None
    
    @staticmethod
    def _extract_ppt_fallback(file_content: bytes, filename: str, start_time: datetime) -> Dict[str, Any]:
        """Fallback extraction for .ppt files when LibreOffice is not available"""
        try:
            # Try using python-pptx anyway (sometimes works with older files)
            logger.info(f"Attempting fallback extraction for {filename}")
            
            try:
                presentation = Presentation(io.BytesIO(file_content))
                return PowerPointExtractor._extract_pptx_content(file_content, filename, start_time)
            except:
                pass
            
            # If that fails, create a basic response with limited info
            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            
            logger.warning(f"Could not extract text from {filename} - file format not fully supported")
            
            extracted_content = {
                "filename": filename,
                "file_size_bytes": len(file_content),
                "slide_count": 0,
                "slides": [],
                "all_text_combined": "",
                "word_count": 0,
                "character_count": 0,
                "processing_time_ms": processing_time,
                "metadata": {
                    "total_slides": 0,
                    "has_content": False,
                    "extraction_timestamp": datetime.now().isoformat(),
                    "extractor_version": "fallback",
                    "file_format": "ppt",
                    "error": "Legacy .ppt format requires LibreOffice for full extraction"
                }
            }
            
            return extracted_content
            
        except Exception as e:
            logger.error(f"Fallback extraction failed: {str(e)}")
            raise
